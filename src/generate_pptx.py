"""
PowerPoint Presentation Generator for Nifty 500 Stock Prediction Project.
Generates an 11-slide widescreen (16:9) executive presentation (.pptx)
using modern corporate styling, dark theme, formatted tables, and KPI callouts.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT_DIR = Path(__file__).resolve().parent.parent
PRESENTATION_DIR = ROOT_DIR / "presentation"
OUTPUT_PPTX = PRESENTATION_DIR / "nifty500_final_presentation.pptx"

# Color Palette: Modern Fintech Dark Theme
BG_COLOR = RGBColor(15, 23, 42)        # Deep Navy / Slate (#0F172A)
CARD_BG = RGBColor(30, 41, 59)        # Card Slate (#1E293B)
PRIMARY = RGBColor(56, 189, 248)      # Electric Cyan (#38BDF8)
SECONDARY = RGBColor(16, 185, 129)    # Emerald (#10B981)
ACCENT = RGBColor(245, 158, 11)       # Amber Gold (#F59E0B)
DANGER = RGBColor(239, 68, 68)        # Red (#EF4444)
TEXT_WHITE = RGBColor(248, 250, 252)  # White (#F8FAFC)
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate Gray (#94A3B8)
BORDER_COLOR = RGBColor(51, 65, 85)   # Border (#334155)


def create_solid_background(slide, color=BG_COLOR):
    """Fills slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_header(slide, title_text: str, category_text: str = "NIFTY 500 PREDICTIVE MODELING · PRD v1.1", slide_num: str = ""):
    """Adds a standard modern header to a slide."""
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    # Category Pill
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = PRIMARY
    p_cat.font.name = "Arial"

    # Title
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.font.name = "Arial"

    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.5), Inches(1.5), Inches(0.5))
        ntf = num_box.text_frame
        np_p = ntf.paragraphs[0]
        np_p.text = slide_num
        np_p.font.size = Pt(11)
        np_p.font.color.rgb = TEXT_MUTED
        np_p.alignment = PP_ALIGN.RIGHT


def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER_COLOR):
    """Creates a rounded rectangle container card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def build_presentation():
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    create_solid_background(s1, BG_COLOR)

    # Accent decorative glow bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(1.2), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # Main Title
    tb_title = s1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.8))
    tf1 = tb_title.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "DATA ANALYTICS INTERN CAPSTONE · PRODUCT REQUIREMENTS DOCUMENT v1.1"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = SECONDARY
    p_badge.font.name = "Arial"

    p_main = tf1.add_paragraph()
    p_main.text = "Nifty 500 Stock Price Prediction System"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.font.name = "Arial"
    p_main.space_before = Pt(8)

    p_sub = tf1.add_paragraph()
    p_sub.text = "End-to-End Analytics Lifecycle: Official NSE Sourcing, BSE Proxy Reconciliation, Econometric Benchmarking & Streamlit Deployment"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.font.name = "Arial"
    p_sub.space_before = Pt(8)

    # Status Pill Card
    add_card(s1, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.5), fill_color=CARD_BG, border_color=BORDER_COLOR)
    tb_meta = s1.shapes.add_textbox(Inches(1.1), Inches(5.25), Inches(11.1), Inches(1.2))
    tf_meta = tb_meta.text_frame
    tf_meta.word_wrap = True

    p_st = tf_meta.paragraphs[0]
    p_st.text = "Project Status: PRD Functionality Implemented · Naive-Baseline Target Not Achieved"
    p_st.font.size = Pt(13)
    p_st.font.bold = True
    p_st.font.color.rgb = ACCENT

    p_dt = tf_meta.add_paragraph()
    p_dt.text = "Dataset: 1,240 Trading Sessions (01 Sep 2021 – 31 Aug 2026) | Primary Benchmark: NSE Nifty 500 | Proxy: BSE 500\nAuthor: Data Analytics Intern | Mentor Review Ready | September 2026"
    p_dt.font.size = Pt(11)
    p_dt.font.color.rgb = TEXT_MUTED
    p_dt.space_before = Pt(4)

    # ==========================================
    # SLIDE 2: Executive Summary & PRD Scope
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    create_solid_background(s2)
    add_slide_header(s2, "Executive Summary & PRD Scope", slide_num="Slide 2 of 11")

    # Card 1: Core Objectives
    add_card(s2, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s2_c1 = s2.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf2_1 = tb_s2_c1.text_frame
    tf2_1.word_wrap = True

    p = tf2_1.paragraphs[0]
    p.text = "🎯 Core Project Objectives (PRD v1.1)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    bullets_s2 = [
        ("Authoritative Data Ingestion:", " 5 years of daily market data (1,240 trading sessions) from the official NSE download with zero web scrapers."),
        ("BSE-500 Cross-Market Proxy:", " Cross-exchange consistency validation across 1,229 common trading sessions."),
        ("Trading Calendar Integrity:", " Strictly preserved active exchange trading days (~250/yr) with 0 synthetic weekend/holiday observations."),
        ("Multi-Family Modeling Suite:", " Benchmarked Persistence Baseline, 5-Day SMA, Statistical ARIMA(1,1,1), Random Forest, XGBoost, and PyTorch LSTM."),
        ("Econometric Rigor:", " Evaluated level RMSE and directional hit rates via formal one-sided Binomial hypothesis testing.")
    ]
    for bold_prefix, text in bullets_s2:
        p = tf2_1.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "• " + bold_prefix
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = text
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: Key Milestones
    add_card(s2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s2_c2 = s2.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf2_2 = tb_s2_c2.text_frame
    tf2_2.word_wrap = True

    p = tf2_2.paragraphs[0]
    p.text = "📊 Key Deliverables & Status"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    milestones = [
        ("Missing Data Benchmark:", " 0.00% missing data post-cleaning (Exceeded PRD target < 2%)."),
        ("Outlier Handling:", " 2 detected return shocks (|Z| > 5) verified as historical macroeconomic events and retained."),
        ("Cross-Validation:", " 5-Fold expanding TimeSeriesSplit cross-validation to assess multi-year stability."),
        ("Performance Target:", " Naive persistence baseline achieved lowest RMSE (209.33) per Martingale law. Target not achieved (honestly reported)."),
        ("Interactive Dashboard:", " Full 8-page Streamlit dashboard with real model-driven multi-step forecasting ($T+1$ to $T+30$).")
    ]
    for bold_prefix, text in milestones:
        p = tf2_2.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "✓ " + bold_prefix
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = text
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: Authoritative Sourcing & BSE Proxy
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    create_solid_background(s3)
    add_slide_header(s3, "Authoritative Data Sourcing & BSE Proxy Architecture", slide_num="Slide 3 of 11")

    # 3 Stat Cards
    stats_s3 = [
        ("1,240 SESSIONS", "Official NSE Historical File", "01 Sep 2021 – 31 Aug 2026", PRIMARY),
        ("0.9999", "Price Correlation", "NSE Nifty 500 vs. BSE 500", SECONDARY),
        ("0.9969", "Daily Return Correlation", "Strong Broad-Market Dynamics", ACCENT)
    ]
    card_w = Inches(3.64)
    for i, (val, title, sub, col) in enumerate(stats_s3):
        x = Inches(0.8 + i * 4.03)
        add_card(s3, x, Inches(1.8), card_w, Inches(1.6))
        tb = s3.shapes.add_textbox(x + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p_val = tf.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(24)
        p_val.font.bold = True
        p_val.font.color.rgb = col
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_s = tf.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(10)
        p_s.font.color.rgb = TEXT_MUTED

    # Detail Card below
    add_card(s3, Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.1))
    tb_det = s3.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(11.1), Inches(2.7))
    tf_det = tb_det.text_frame
    tf_det.word_wrap = True

    p = tf_det.paragraphs[0]
    p.text = "🏛️ Sourcing Architecture & Role Clarification"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    det_bullets = [
        ("Primary Authoritative Source:", " Sourced directly from official NSE historical index download (NIFTY_500_Historical_PR_01-09-2021 to 31-08-2026.csv). Sourced with full schema: Date, Open, High, Low, Close, Volume, Adjusted Close. Zero third-party web scrapers."),
        ("BSE-500 Role Clarification:", " BSE-500 is utilized strictly as a cross-exchange broad-market proxy for co-movement and reconciliation per PRD Section 5. It is NOT an interchangeable clone of Nifty 500."),
        ("Reconciliation Findings:", " Evaluated across 1,229 common trading days. The near-perfect correlation (0.9999 price, 0.9969 return) confirms that both benchmarks reflect identical macro equity drivers across India's two premier exchanges.")
    ]
    for b_prefix, b_text in det_bullets:
        p = tf_det.add_paragraph()
        p.space_before = Pt(8)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "• " + b_prefix
        r1.font.bold = True
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = b_text
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 4: Trading Calendar & Outlier Handling
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    create_solid_background(s4)
    add_slide_header(s4, "Trading Calendar Integrity & Outlier Event Validation", slide_num="Slide 4 of 11")

    # Left: Outlier Table Card
    add_card(s4, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.0))
    tb_ot = s4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(6.9), Inches(4.5))
    tf_ot = tb_ot.text_frame
    tf_ot.word_wrap = True

    p = tf_ot.paragraphs[0]
    p.text = "🚨 Verified Outlier Events (|Z-Score| > 5)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    outliers_data = [
        ("Event 1: Russia-Ukraine Crisis Outbreak", "Date: 2022-02-24 | Close: ₹13,775.70 | Daily Return: -5.04% (Z = -5.59)\nContext: Global equity market crash upon geopolitical conflict outbreak. Verified against official NSE exchange records."),
        ("Event 2: 2024 Indian General Election Results", "Date: 2024-06-04 | Close: ₹20,323.85 | Daily Return: -6.76% (Z = -7.48)\nContext: Historic intraday market shock following unexpected election margin narrowness. Verified against official NSE exchange records.")
    ]
    for title, desc in outliers_data:
        p = tf_ot.add_paragraph()
        p.space_before = Pt(12)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.text = title

        p2 = tf_ot.add_paragraph()
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.text = desc

    # Right: Retention Rationale Card
    add_card(s4, Inches(8.6), Inches(1.8), Inches(3.9), Inches(5.0))
    tb_rr = s4.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(3.5), Inches(4.5))
    tf_rr = tb_rr.text_frame
    tf_rr.word_wrap = True

    p = tf_rr.paragraphs[0]
    p.text = "💡 Econometric Retention Rationale"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    reasons = [
        ("No Recording Errors:", " Neither date is corrupt or misprinted. Both are verified, historic macroeconomic shocks."),
        ("Censorship Bias Prevention:", " Deleting or capping genuine market shocks artificially suppresses return variance and understates downside risk."),
        ("Fat-Tail Reality:", " Preserving verified shocks ensures models learn genuine leptokurtic equity dynamics."),
        ("Calendar Discipline:", " Zero synthetic observations manufactured for weekends or exchange holidays.")
    ]
    for b_title, b_desc in reasons:
        p = tf_rr.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(10)
        r1 = p.add_run()
        r1.text = "✓ " + b_title
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 5: Feature Engineering Architecture
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    create_solid_background(s5)
    add_slide_header(s5, "Feature Engineering Architecture (15+ Technical Signals)", slide_num="Slide 5 of 11")

    # 4 Category Cards
    feat_cats = [
        ("Trend Indicators", "SMA 20, SMA 50, SMA 200\nEMA 20, EMA 50\nPrice-to-MA Distance Ratios", PRIMARY),
        ("Momentum Oscillators", "Relative Strength Index (RSI 14)\nMACD Line (12/26-day)\nMACD Signal Line (9-day)\nMACD Histogram", SECONDARY),
        ("Volatility Measures", "Bollinger Bands (Upper/Lower/Width)\nBollinger %B Indicator\n20-Day Annualized Volatility\n50-Day Annualized Volatility", ACCENT),
        ("Volume & Lags", "Volume 20-Day SMA\nVolume Ratio to Moving Average\nPrice Lags (t-1, t-2, t-5)\nReturn Lags (1d, 5d, 20d)", RGBColor(168, 85, 247))
    ]
    card_w2 = Inches(2.7)
    for i, (title, content, col) in enumerate(feat_cats):
        x = Inches(0.8 + i * 3.0)
        add_card(s5, x, Inches(1.8), card_w2, Inches(3.4))
        tb = s5.shapes.add_textbox(x + Inches(0.2), Inches(2.0), card_w2 - Inches(0.4), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p_c = tf.add_paragraph()
        p_c.text = content
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = TEXT_MUTED
        p_c.space_before = Pt(10)

    # Bottom Banner: Leakage Prevention
    add_card(s5, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3))
    tb_lk = s5.shapes.add_textbox(Inches(1.1), Inches(5.65), Inches(11.1), Inches(1.0))
    tf_lk = tb_lk.text_frame
    tf_lk.word_wrap = True
    p = tf_lk.paragraphs[0]
    p.text = "🛡️ Zero Lookahead Bias Discipline"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p_sub = tf_lk.add_paragraph()
    p_sub.text = "Warmup period of 200 trading days filtered out (1,040 feature rows retained). Every indicator is strictly computed using market information available on day t to predict Close(t+1). Scalers fitted exclusively on the training partition."
    p_sub.font.size = Pt(10.5)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(4)

    # ==========================================
    # SLIDE 6: Chronological Validation & Cross-Validation
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    create_solid_background(s6)
    add_slide_header(s6, "Validation Methodology & Time-Series Cross-Validation", slide_num="Slide 6 of 11")

    # Left: Holdout Split Card
    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s6_1 = s6.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf6_1 = tb_s6_1.text_frame
    tf6_1.word_wrap = True
    p = tf6_1.paragraphs[0]
    p.text = "⏱️ Chronological 80/20 Holdout Split"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    split_items = [
        ("Strict Temporal Split:", " No random shuffling to prevent future data leakage into training."),
        ("Training Window:", " 832 trading sessions (21 Jun 2022 to 27 Oct 2025). Used for indicator baselining and model fitting."),
        ("Out-of-Sample Test Window:", " 208 trading sessions (28 Oct 2025 to 28 Aug 2026; ~10 calendar months). Evaluated strictly out-of-sample."),
        ("Leak-Proof ARIMA Walk-Forward:", " Iterative 1-step rolling extend (model.extend()) without re-smoothing future test observations.")
    ]
    for h, b in split_items:
        p = tf6_1.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # Right: 5-Fold TimeSeriesSplit CV Card
    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s6_2 = s6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf6_2 = tb_s6_2.text_frame
    tf6_2.word_wrap = True
    p = tf6_2.paragraphs[0]
    p.text = "🔄 5-Fold TimeSeriesSplit Cross-Validation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    cv_desc = [
        ("PRD Section 12 Requirement:", " Mitigate overfitting risk by evaluating stability across expanding historical horizons."),
        ("Random Forest CV RMSE:", " 1,094.54 (±779.47) | CV MAE: 849.42"),
        ("XGBoost Regressor CV RMSE:", " 1,067.78 (±828.08) | CV MAE: 832.33"),
        ("Econometric Finding:", " As Nifty 500 expanded from ₹13,000 to ₹23,000 across 2022–2025, error variance scales proportionally with index level, highlighting multi-regime dynamics.")
    ]
    for h, b in cv_desc:
        p = tf6_2.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "✓ " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: Benchmark Scorecard
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    create_solid_background(s7)
    add_slide_header(s7, "Model Performance Benchmark Scorecard (N = 208 Test Days)", slide_num="Slide 7 of 11")

    # Table Shape
    rows = 7
    cols = 6
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.7)
    height = Inches(3.6)

    table_shape = s7.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.6)
    table.columns[4].width = Inches(1.8)
    table.columns[5].width = Inches(1.7)

    headers = ["Model Architecture", "Family", "RMSE (Pts)", "MAE (Pts)", "MAPE (%)", "Directional Acc"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT

    scorecard_data = [
        ("Naive Persistence Baseline", "Benchmark", "209.33", "151.88", "0.669%", "N/A", True),
        ("Random Forest Regressor", "Classical ML", "228.48", "169.05", "0.742%", "48.56%", False),
        ("XGBoost Regressor", "Classical ML", "239.19", "181.36", "0.795%", "51.92%", False),
        ("5-Day Moving Average (SMA)", "Benchmark", "288.86", "218.20", "0.957%", "51.44%", False),
        ("ARIMA(1, 1, 1) Walk-Forward", "Statistical", "291.08", "216.41", "0.950%", "51.92%", False),
        ("PyTorch LSTM Network", "Deep Learning", "563.46", "425.56", "1.877%", "56.73%", False)
    ]

    for i, (m_name, fam, rmse, mae, mape, d_acc, is_best) in enumerate(scorecard_data, 1):
        row_vals = [m_name, fam, rmse, mae, mape, d_acc]
        for j, val in enumerate(row_vals):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(16, 45, 38) if is_best else CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10.5)
            p.font.name = "Arial"
            p.font.bold = is_best
            p.font.color.rgb = SECONDARY if is_best else (TEXT_WHITE if j == 0 else TEXT_MUTED)
            p.alignment = PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT

    # Callout Banner below Table
    add_card(s7, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3))
    tb_sc = s7.shapes.add_textbox(Inches(1.1), Inches(5.75), Inches(11.1), Inches(1.0))
    tf_sc = tb_sc.text_frame
    tf_sc.word_wrap = True

    p = tf_sc.paragraphs[0]
    p.text = "⚖️ Critical Finding: Naive Persistence Baseline Achieves Lowest Level-Price RMSE (209.33)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p_sc = tf_sc.add_paragraph()
    p_sc.text = "Status: PRD functionality implemented; naive-baseline performance target not achieved. Under the Martingale property of asset prices (E[Pt+1 | Ft] = Pt), today's price is the minimum-variance quadratic estimator of tomorrow's price level."
    p_sc.font.size = Pt(10)
    p_sc.font.color.rgb = TEXT_MUTED
    p_sc.space_before = Pt(3)

    # ==========================================
    # SLIDE 8: Statistical Significance & Econometrics
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    create_solid_background(s8)
    add_slide_header(s8, "Statistical Hypothesis Testing & Directional Realities", slide_num="Slide 8 of 11")

    # Card 1: Binomial Hypothesis Testing
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s8_1 = s8.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf8_1 = tb_s8_1.text_frame
    tf8_1.word_wrap = True
    p = tf8_1.paragraphs[0]
    p.text = "🎲 Binomial Directional Tests (H0: p = 0.50)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    b_tests = [
        ("Naive Persistence Baseline:", " Directional Accuracy = N/A. Since Pt+1 = Pt, it predicts zero price movement, not a directional signal."),
        ("XGBoost Regressor:", " 108 / 208 days (51.92%) | 1-Sided p-value: 0.3138 | 95% Wilson CI: [45.16%, 58.62%]. Fail to reject H0."),
        ("PyTorch LSTM Network:", " 118 / 208 days (56.73%) | 1-Sided p-value: 0.0305 | 95% Wilson CI: [49.94%, 63.28%]. Statistically significant directional momentum signal, but exhibits severe level-price drift (RMSE 563.46).")
    ]
    for h, b in b_tests:
        p = tf8_1.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(10.5)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: Deep Learning Non-Stationarity
    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s8_2 = s8.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf8_2 = tb_s8_2.text_frame
    tf8_2.word_wrap = True
    p = tf8_2.paragraphs[0]
    p.text = "🧠 The Deep Learning Trade-Off"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    dl_insights = [
        ("Why LSTM Struggles on Levels:", " Raw equity price series are non-stationary with stochastic drift. Neural networks trained on MinMax-scaled price levels lack explicit local mean-reversion anchors, causing predictions to lag turning points."),
        ("Sequence Momentum:", " The LSTM network successfully captures sequential regime momentum (56.73% hit rate), confirming multi-day temporal memory."),
        ("Institutional Practice:", " In quantitative hedge funds, deep learning models are rarely trained on raw price levels; instead, they are formulated on stationary returns or residual alpha.")
    ]
    for h, b in dl_insights:
        p = tf8_2.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(10.5)
        r1 = p.add_run()
        r1.text = "✓ " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 9: Real Model-Driven Future Forecasting
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    create_solid_background(s9)
    add_slide_header(s9, "Future Horizon Target Forecaster (Real Model Rollouts)", slide_num="Slide 9 of 11")

    # Card 1: Recursive ML Engine
    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s9_1 = s9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf9_1 = tb_s9_1.text_frame
    tf9_1.word_wrap = True
    p = tf9_1.paragraphs[0]
    p.text = "🔮 Recursive Machine Learning Rollouts"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    fc_details = [
        ("Trained Model Driven:", " Unlike static fixed mathematical drift, future target forecasts are generated by actual trained models (Recursive XGBoost, Random Forest, ARIMA)."),
        ("Iterative Indicator Recalculation:", " At each forward step t+1...t+H, the predicted price is appended to the series and all 15+ technical indicators (SMA, EMA, RSI, MACD, BB) are dynamically updated."),
        ("Compounding Uncertainty Bands:", " Confidence intervals expand realistically with forecast horizon length: Bounds = Forecast ± (z * RMSE * sqrt(h)).")
    ]
    for h, b in fc_details:
        p = tf9_1.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(10.5)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: Two-Tier Horizon Design
    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s9_2 = s9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf9_2 = tb_s9_2.text_frame
    tf9_2.word_wrap = True
    p = tf9_2.paragraphs[0]
    p.text = "📐 Two-Tier Horizon Architecture"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    horizons = [
        ("Tier 1 — Backtesting Target (T+1):", " 1-step ahead daily closing price evaluation for rigorous sequential testing against persistence baselines with zero lookahead bias."),
        ("Tier 2 — Future Scenario Horizon (T+1 to T+30):", " User-selectable multi-day forward projections in the Streamlit application for tactical scenario planning."),
        ("Statistical ARIMA Rollout:", " Generated via statsmodels get_forecast(steps=H) with covariance-matrix derived confidence intervals.")
    ]
    for h, b in horizons:
        p = tf9_2.add_paragraph()
        p.space_before = Pt(10)
        p.font.size = Pt(10.5)
        r1 = p.add_run()
        r1.text = "✓ " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 10: Interactive Streamlit Dashboard
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    create_solid_background(s10)
    add_slide_header(s10, "Interactive Streamlit Dashboard Architecture", slide_num="Slide 10 of 11")

    # 4 Feature Tiles
    tiles = [
        ("Executive Overview & KPIs", "Real-time index metrics, 52-week boundaries, 5-year interactive Plotly price chart with 50/200 SMA overlays.", PRIMARY),
        ("Technical Analysis & EDA", "Candlestick charts, Bollinger Bands, RSI, MACD, volatility clustering, and verified outlier shock table (|Z| > 5).", SECONDARY),
        ("Cross-Exchange Reconciliation", "NSE vs. BSE correlation analysis (r = 0.9999 price, r = 0.9969 return) and trading calendar alignment.", ACCENT),
        ("Model Scorecard & Forecaster", "Full benchmark leaderboard, 5-Fold cross-validation display, and interactive T+1 to T+30 recursive forecasting.", RGBColor(168, 85, 247))
    ]
    for i, (t_title, t_desc, col) in enumerate(tiles):
        row_i = i // 2
        col_i = i % 2
        x = Inches(0.8 + col_i * 6.0)
        y = Inches(1.8 + row_i * 2.6)
        add_card(s10, x, y, Inches(5.7), Inches(2.3))
        tb = s10.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p_c = tf.add_paragraph()
        p_c.text = t_desc
        p_c.font.size = Pt(10.5)
        p_c.font.color.rgb = TEXT_MUTED
        p_c.space_before = Pt(6)

    # ==========================================
    # SLIDE 11: Limitations, Disclaimers & Next Steps
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    create_solid_background(s11)
    add_slide_header(s11, "Project Limitations, Academic Disclaimer & Next Steps", slide_num="Slide 11 of 11")

    # Left: Limitations & Disclaimers
    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s11_1 = s11.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf11_1 = tb_s11_1.text_frame
    tf11_1.word_wrap = True
    p = tf11_1.paragraphs[0]
    p.text = "⚠️ Limitations & Academic Disclaimers"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DANGER

    limits = [
        ("Unmodeled Execution Frictions:", " Bid-ask spreads, broker transaction fees, and Securities Transaction Tax (STT) are unmodeled."),
        ("Structural Regime Shifts:", " Major macroeconomic regime transitions cannot be predicted purely from historical price patterns."),
        ("Academic Portfolio Project:", " Developed strictly as an educational intern research project. It is NOT financial advice, nor is it intended for live automated trading or real capital deployment.")
    ]
    for h, b in limits:
        p = tf11_1.add_paragraph()
        p.space_before = Pt(12)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "• " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # Right: Recommended Next Steps
    add_card(s11, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_s11_2 = s11.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf11_2 = tb_s11_2.text_frame
    tf11_2.word_wrap = True
    p = tf11_2.paragraphs[0]
    p.text = "🚀 Recommended Quantitative Next Steps"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    next_steps = [
        ("Stationary Returns Formulation:", " Train deep learning architectures (LSTM, Temporal Fusion Transformers) on stationary log returns rather than raw price levels."),
        ("Macroeconomic Feature Ingestion:", " Incorporate RBI repo rates, USD/INR exchange rates, Brent crude oil futures, and India VIX."),
        ("Hierarchical Sector Models:", " Forecast broad index movements by modeling top constituent sectors (Nifty Bank, Nifty IT, Nifty Auto).")
    ]
    for h, b in next_steps:
        p = tf11_2.add_paragraph()
        p.space_before = Pt(12)
        p.font.size = Pt(11)
        r1 = p.add_run()
        r1.text = "✓ " + h
        r1.font.bold = True
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = TEXT_MUTED

    # Save presentation
    prs.save(OUTPUT_PPTX)
    print(f"[PPTX GENERATOR] Successfully created executive presentation ({len(prs.slides)} slides) at: {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_presentation()
